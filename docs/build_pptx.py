#!/usr/bin/env python3
"""
build_pptx.py - Sinh bài thuyết trình (.pptx) cho Đề tài 42 - Networking.

Dùng: python build_pptx.py
Yêu cầu: pip install python-pptx

~21 slide: tổng quan -> cơ sở lý thuyết networking -> 5 bài kernel module
(kèm trích code & phân tích) -> web demo -> kết luận. Part 1 & 2 chỉ điểm qua.
Các chỗ cần ảnh được chèn placeholder kèm gợi ý nội dung ảnh.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "ThuyetTrinh_DeTai42_Networking.pptx")

# ----- Bảng màu -----
NAVY   = RGBColor(0x1F, 0x47, 0x7C)
BLUE   = RGBColor(0x2E, 0x74, 0xB5)
LIGHT  = RGBColor(0xEA, 0xF1, 0xF8)
CODEBG = RGBColor(0x1E, 0x29, 0x3B)
CODEFG = RGBColor(0x9C, 0xDC, 0xFE)
GRAY   = RGBColor(0x55, 0x55, 0x55)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PHBG   = RGBColor(0xF3, 0xF6, 0xFA)
PHBD   = RGBColor(0xB5, 0xC7, 0xDD)
MONO   = "Consolas"
SANS   = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(tf_or_para, text, size, color, bold=False, font=SANS, align=None):
    p = tf_or_para
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    if align is not None:
        p.alignment = align


def add_box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, round=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
        l, t, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def title_bar(slide, title, kicker=None):
    add_box(slide, 0, 0, SW, Inches(1.15), fill=NAVY)
    add_box(slide, 0, Inches(1.15), SW, Pt(4), fill=BLUE)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), SW - Inches(1), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    if kicker:
        _set(tf.paragraphs[0], kicker, 12, RGBColor(0xBF, 0xD6, 0xF0), bold=True)
        p = tf.add_paragraph(); _set(p, title, 26, WHITE, bold=True)
    else:
        _set(tf.paragraphs[0], title, 28, WHITE, bold=True)
        tf.paragraphs[0].space_before = Pt(8)


def footer(slide, idx):
    tb = slide.shapes.add_textbox(Inches(0.4), SH - Inches(0.42),
                                  SW - Inches(0.8), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    _set(p, "Đề tài 42 — Lập trình nhân Linux: Networking", 9, GRAY)
    n = slide.shapes.add_textbox(SW - Inches(1.2), SH - Inches(0.42),
                                 Inches(0.9), Inches(0.3))
    pn = n.text_frame.paragraphs[0]
    _set(pn, str(idx), 10, GRAY, align=PP_ALIGN.RIGHT)


_slide_no = 0


def new_slide(title=None, kicker=None, count=True):
    global _slide_no
    s = prs.slides.add_slide(BLANK)
    add_box(s, 0, 0, SW, SH, fill=WHITE)
    if title:
        title_bar(s, title, kicker)
    if count:
        _slide_no += 1
        footer(s, _slide_no)
    return s


def bullets(slide, items, left=Inches(0.7), top=Inches(1.5),
            width=None, size=18, gap=10):
    width = width or (SW - Inches(1.4))
    tb = slide.shapes.add_textbox(left, top, width, SH - top - Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0
        txt = it
        if isinstance(it, tuple):
            # tuple = sub-bullet (level 1); phần tử đầu chỉ là đánh dấu
            lvl, txt = 1, it[1]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet = "•  " if lvl == 0 else "–  "
        _set(p, bullet + txt, size - lvl * 2,
             NAVY if lvl == 0 else GRAY, bold=(lvl == 0))
        p.level = lvl
        p.space_after = Pt(gap)
    return tb


def code_box(slide, code, left=Inches(0.6), top=Inches(1.45),
             width=None, height=None, size=12, caption=None):
    width = width or (SW - Inches(1.2))
    height = height or Inches(4.3)
    box = add_box(slide, left, top, width, height, fill=CODEBG, round=True)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP          # canh trên (không giữa)
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.1)
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set(p, line if line else " ", size, CODEFG, font=MONO,
             align=PP_ALIGN.LEFT)             # ép căn trái mọi dòng
        p.line_spacing = 1.0; p.space_after = Pt(0)
    if caption:
        cb = slide.shapes.add_textbox(left, top + height + Inches(0.05),
                                      width, Inches(0.4))
        _set(cb.text_frame.paragraphs[0], caption, 12, GRAY, font=MONO)
    return box


def img_ph(slide, suggestion, left=None, top=Inches(1.5),
           width=Inches(5.4), height=Inches(4.6)):
    left = left if left is not None else (SW - width - Inches(0.6))
    box = add_box(slide, left, top, width, height, fill=PHBG,
                  line=PHBD, line_w=1.5, round=True)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    _set(tf.paragraphs[0], "🖼  CHÈN ẢNH", 16, BLUE, bold=True, align=PP_ALIGN.CENTER)
    p = tf.add_paragraph()
    _set(p, suggestion, 12, GRAY, align=PP_ALIGN.CENTER)
    p.space_before = Pt(8)
    return box


def purpose(slide, text, top=Inches(1.32), width=None):
    """Dải callout 'Mục đích' đặt ngay dưới thanh tiêu đề."""
    width = width or (SW - Inches(1.2))
    box = add_box(slide, Inches(0.6), top, width, Inches(0.7),
                  fill=LIGHT, line=BLUE, line_w=1.0, round=True)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "🎯 Mục đích:  "
    r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = NAVY
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(15); r2.font.color.rgb = GRAY
    return box


# =====================================================================
# 1. TITLE
# =====================================================================
s = new_slide(count=False)
add_box(s, 0, 0, SW, SH, fill=NAVY)
add_box(s, 0, Inches(3.4), SW, Pt(3), fill=BLUE)
tb = s.shapes.add_textbox(Inches(1), Inches(1.0), SW - Inches(2), Inches(1))
_set(tb.text_frame.paragraphs[0], "HỌC VIỆN KỸ THUẬT MẬT MÃ — KHOA CNTT",
     16, RGBColor(0xBF, 0xD6, 0xF0), bold=True, align=PP_ALIGN.CENTER)
tb = s.shapes.add_textbox(Inches(1), Inches(2.0), SW - Inches(2), Inches(1.3))
_set(tb.text_frame.paragraphs[0], "LẬP TRÌNH NHÂN LINUX", 30, WHITE,
     bold=True, align=PP_ALIGN.CENTER)
tb = s.shapes.add_textbox(Inches(1), Inches(3.7), SW - Inches(2), Inches(1.6))
tf = tb.text_frame; tf.word_wrap = True
_set(tf.paragraphs[0], "Đề tài 42 — Networking", 24,
     RGBColor(0x8E, 0xC6, 0xFF), bold=True, align=PP_ALIGN.CENTER)
p = tf.add_paragraph()
_set(p, "Xây dựng mô-đun nhân xử lý mạng & tích hợp vào hệ thống",
     16, WHITE, align=PP_ALIGN.CENTER)
tb = s.shapes.add_textbox(Inches(1), Inches(5.7), SW - Inches(2), Inches(1.2))
tf = tb.text_frame
_set(tf.paragraphs[0], "Nhóm thực hiện: Ngô Hồng Phong · Ngô Minh Cường · Đỗ Minh Thuần",
     14, WHITE, align=PP_ALIGN.CENTER)
p = tf.add_paragraph(); _set(p, "Hà Nội, 2026", 13,
                             RGBColor(0xBF, 0xD6, 0xF0), align=PP_ALIGN.CENTER)

# =====================================================================
# 2. NỘI DUNG
# =====================================================================
s = new_slide("Nội dung trình bày")
agenda = [
    ("1", "Tổng quan đề tài", "Mục tiêu & phạm vi"),
    ("2", "Kiến thức nền", "Nhân Linux · kernel space · kernel module"),
    ("3", "Mạng trong nhân Linux", "Network stack · sk_buff · socket · Netfilter"),
    ("4", "Triển khai 5 phần kernel module", "Giám sát · Lọc IP · TCP · UDP"),
    ("5", "Kết quả thực nghiệm", "Demo & đánh giá từng phần"),
    ("6", "Kết luận", "Tổng kết & hướng phát triển"),
]
y = Inches(1.55)
row_h = Inches(0.84)
for i, (num, title, sub) in enumerate(agenda):
    add_box(s, Inches(0.9), y, Inches(11.5), Inches(0.72),
            fill=(LIGHT if i % 2 == 0 else WHITE),
            line=PHBD, line_w=0.75, round=True)
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.08),
                               y + Inches(0.09), Inches(0.54), Inches(0.54))
    badge.fill.solid(); badge.fill.fore_color.rgb = NAVY
    badge.line.fill.background(); badge.shadow.inherit = False
    _set(badge.text_frame.paragraphs[0], num, 20, WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    tb = s.shapes.add_textbox(Inches(1.95), y + Inches(0.02),
                              Inches(10.2), Inches(0.68))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title + "    "
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = NAVY
    r2 = p.add_run(); r2.text = sub
    r2.font.size = Pt(13); r2.font.color.rgb = GRAY; r2.font.italic = True
    y = y + row_h

# =====================================================================
# 3. TỔNG QUAN ĐỀ TÀI
# =====================================================================
s = new_slide("Tổng quan đề tài", kicker="GIỚI THIỆU")
bullets(s, [
    "Đề tài gồm 3 nội dung, trọng tâm là lập trình mạng ở mức nhân:",
    ("", "Shell: quản lý file, lập lịch, thời gian, cài/gỡ phần mềm"),
    ("", "User-space: quản lý tiến trình, file, socket/network"),
    ("", "Kernel module Networking (TRỌNG TÂM)"),
    "Mục tiêu phần Networking:",
    ("", "Hiểu network stack & đường đi của gói tin trong nhân"),
    ("", "Bắt/giám sát & lọc gói tin bằng Netfilter hook"),
    ("", "Lập trình socket trong không gian nhân (TCP/UDP)"),
    ("", "Giao tiếp user ↔ kernel qua ioctl"),
], size=18, gap=9)

# =====================================================================
# KIẾN THỨC NỀN (bổ sung) — Nhân Linux & hệ thống mạng
# =====================================================================
s = new_slide("Nhân Linux (Linux kernel) là gì?", kicker="KIẾN THỨC NỀN")
bullets(s, [
    "Lõi của hệ điều hành — trung gian giữa phần cứng & ứng dụng.",
    "Quản lý: tiến trình, bộ nhớ, file, thiết bị và mạng.",
    "Kiến trúc nguyên khối (monolithic): dịch vụ chạy chung không gian nhân.",
    ("", "Hiệu năng cao nhưng 1 lỗi nhỏ có thể sập hệ thống."),
    "Viết chủ yếu bằng C, mã nguồn mở.",
    "Ứng dụng truy cập phần cứng gián tiếp qua lời gọi hệ thống (syscall).",
], left=Inches(0.6), width=Inches(6.3), size=16, gap=11)
code_box(s,
    "/* Mã chạy trong nhân dùng các kiểu\n"
    "   dữ liệu & API riêng của kernel.\n"
    "   VD chữ ký hàm xử lý gói tin: */\n"
    "\n"
    "typedef unsigned int nf_hookfn(\n"
    "        void *priv,\n"
    "        struct sk_buff *skb,\n"
    "        const struct nf_hook_state *state);",
    left=Inches(7.1), width=Inches(5.6), top=Inches(1.7),
    height=Inches(3.3), size=12.5,
    caption="Nguồn: Linux Kernel Labs — Networking")

s = new_slide("Kernel space vs User space", kicker="KIẾN THỨC NỀN")
bullets(s, [
    "Bộ nhớ & quyền thực thi chia làm 2 vùng tách biệt:",
    ("", "User space (ring 3): ứng dụng thường, quyền hạn chế, bị cô lập."),
    ("", "Kernel space (ring 0): nhân, toàn quyền với phần cứng."),
    "Cầu nối giữa hai vùng là system call (open, read, socket...).",
    "Code chạy ở kernel space:",
    ("", "Truy cập trực tiếp bộ nhớ/phần cứng, không được bảo vệ như user."),
    ("", "Lỗi → kernel panic. ⇒ Phải test trong máy ảo, snapshot trước."),
], left=Inches(0.6), width=Inches(6.6), size=17, gap=10)
img_ph(s, "Sơ đồ 2 vùng User space / Kernel space và syscall làm cầu nối "
          "(ring 3 ↔ ring 0).",
       left=Inches(7.4), width=Inches(5.3), top=Inches(1.5), height=Inches(4.6))

s = new_slide("Kernel module (LKM)", kicker="KIẾN THỨC NỀN")
bullets(s, [
    "Loadable Kernel Module: mã nạp/gỡ vào nhân lúc chạy.",
    ("", "Không cần biên dịch & khởi động lại toàn bộ nhân."),
    "Dùng cho driver, hệ thống file, mô-đun mạng (đề tài này).",
    "Lệnh quản lý:",
    ("", "insmod / rmmod — nạp / gỡ (.ko)"),
    ("", "lsmod · modinfo · dmesg (xem log nhân)"),
    "Mỗi module có hàm init (nạp) & exit (gỡ).",
], left=Inches(0.6), width=Inches(6.3), size=16, gap=11)
code_box(s,
    "#include <linux/module.h>\n"
    "#include <linux/kernel.h>\n"
    "\n"
    "static int __init my_init(void) {\n"
    "    pr_info(\"module: loaded\\n\");\n"
    "    return 0;\n"
    "}\n"
    "static void __exit my_exit(void) {\n"
    "    pr_info(\"module: unloaded\\n\");\n"
    "}\n"
    "\n"
    "module_init(my_init);\n"
    "module_exit(my_exit);\n"
    "MODULE_LICENSE(\"GPL\");",
    left=Inches(7.1), width=Inches(5.6), top=Inches(1.5),
    height=Inches(4.7), size=12.5,
    caption="Khung cơ bản của một kernel module")

s = new_slide("Hệ thống mạng trong Linux", kicker="KIẾN THỨC NỀN")
bullets(s, [
    "Toàn bộ truyền thông mạng do nhân điều phối — đường đi gói tin nhận:",
    ("", "Card mạng (NIC) → driver → ngắt/softirq (NAPI)"),
    ("", "→ tầng IP → tầng TCP/UDP → hàng đợi socket"),
    ("", "→ ứng dụng user-space đọc qua read()/recv()"),
    "Mỗi tầng bọc/gỡ header tương ứng (Ethernet, IP, TCP/UDP).",
    "Linux nổi tiếng vì stack mạng mạnh ⇒ được dùng nhiều cho router, firewall, server.",
], left=Inches(0.6), width=Inches(6.6), size=17, gap=10)
img_ph(s, "Sơ đồ luồng gói tin từ NIC qua các tầng nhân lên ứng dụng "
          "(NIC → driver → IP → TCP/UDP → socket → app).",
       left=Inches(7.4), width=Inches(5.3), top=Inches(1.5), height=Inches(4.6))

s = new_slide("Netfilter & iptables — bức tranh lớn", kicker="KIẾN THỨC NỀN")
bullets(s, [
    "Netfilter là khung lọc gói tích hợp sẵn trong nhân.",
    "iptables / nftables là công cụ user-space cấu hình Netfilter:",
    ("", "tables (filter, nat, mangle) chứa các chain"),
    ("", "mỗi chain gắn vào 1 trong 5 hook"),
    "Tường lửa, NAT, chia sẻ Internet đều dựa trên cơ chế này.",
    "Mô-đun của nhóm đăng ký hook ở cùng tầng ⇒ firewall thu nhỏ.",
], left=Inches(0.6), width=Inches(6.3), size=16, gap=12)
code_box(s,
    "/* Mẫu đăng ký hook */\n"
    "static struct nf_hook_ops my_nfho = {\n"
    "    .hook     = my_nf_hookfn,\n"
    "    .hooknum  = NF_INET_LOCAL_OUT,\n"
    "    .pf       = PF_INET,\n"
    "    .priority = NF_IP_PRI_FIRST\n"
    "};\n"
    "\n"
    "int __init my_hook_init(void) {\n"
    "    return nf_register_net_hook(\n"
    "               &init_net, &my_nfho);\n"
    "}",
    left=Inches(7.1), width=Inches(5.6), top=Inches(1.5),
    height=Inches(4.2), size=12.5,
    caption="Nguồn: Linux Kernel Labs — Networking")

# =====================================================================
# 4. NETWORK STACK
# =====================================================================
s = new_slide("Network stack trong nhân Linux", kicker="CƠ SỞ LÝ THUYẾT")
bullets(s, [
    "Mọi gói tin vào/ra đều đi qua chuỗi xử lý trong nhân.",
    "Phân tầng theo mô hình TCP/IP:",
    ("", "Link → IP → TCP/UDP → Socket → ứng dụng user-space"),
    "Kể cả lưu lượng localhost cũng qua nhân.",
    "→ Can thiệp ở mức nhân = quan sát/điều khiển toàn diện.",
], left=Inches(0.6), width=Inches(6.5), size=18, gap=12)
img_ph(s, "Sơ đồ network stack / đường đi gói tin qua các tầng "
          "(có thể lấy từ trang Linux Kernel Labs — Networking).",
       left=Inches(7.3), width=Inches(5.4), top=Inches(1.5), height=Inches(4.5))

# =====================================================================
# 6. sk_buff
# =====================================================================
s = new_slide("struct sk_buff — gói tin trong nhân", kicker="CƠ SỞ LÝ THUYẾT")
bullets(s, [
    "sk_buff (socket buffer): cấu trúc trung tâm biểu diễn 1 gói tin.",
    "Chứa dữ liệu gói + con trỏ tới từng tầng header.",
    "Hàm trợ giúp truy cập header:",
], left=Inches(0.6), width=Inches(12), size=18, gap=10)
code_box(s, "ip_hdr(skb)    /* struct iphdr  — saddr, daddr, protocol */\n"
            "tcp_hdr(skb)   /* struct tcphdr — source, dest, syn/ack/fin */\n"
            "udp_hdr(skb)   /* struct udphdr — source, dest */",
         top=Inches(3.5), height=Inches(1.7), size=15)

# =====================================================================
# 7. socket & sock
# =====================================================================
s = new_slide("struct socket & struct sock", kicker="CƠ SỞ LÝ THUYẾT")
bullets(s, [
    "struct socket — lớp trừu tượng BSD socket (phía giao tiếp):",
    ("", "ops  → các thao tác theo giao thức (bind/listen/accept...)"),
    ("", "sk   → liên kết tới socket INET nội bộ"),
    "struct sock — biểu diễn nội bộ của socket INET:",
    ("", "Lưu trạng thái kết nối, thông tin giao thức"),
    ("", "Callback khi có sự kiện (dữ liệu đến, đổi trạng thái)"),
], size=18, gap=12)

# =====================================================================
# 8. NETFILTER
# =====================================================================
s = new_slide("Netfilter & các điểm hook", kicker="CƠ SỞ LÝ THUYẾT")
bullets(s, [
    "Framework đăng ký callback (hook) tại các điểm cố định trên đường đi gói tin.",
    "Là nền tảng của iptables/nftables. Đăng ký qua nf_register_net_hook().",
], left=Inches(0.6), width=Inches(12), size=17, gap=8)
# bảng hook
rows = [
    ("Hook", "Vị trí", "Dùng để"),
    ("PRE_ROUTING", "Gói vừa vào, trước định tuyến", "Lọc/giám sát gói ĐẾN"),
    ("LOCAL_IN", "Gói tới tiến trình cục bộ", "Firewall inbound"),
    ("LOCAL_OUT", "Gói do máy tạo ra, đi ra", "Lọc/giám sát gói ĐI"),
    ("POST_ROUTING", "Trước khi rời máy", "NAT nguồn"),
]
tbl_shape = s.shapes.add_table(len(rows), 3, Inches(0.6), Inches(2.5),
                               Inches(9.2), Inches(2.6)).table
tbl_shape.columns[0].width = Inches(2.4)
tbl_shape.columns[1].width = Inches(4.2)
tbl_shape.columns[2].width = Inches(2.6)
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        c = tbl_shape.cell(i, j)
        c.text = val
        pr = c.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(13)
        pr.runs[0].font.name = SANS
        if i == 0:
            pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = WHITE
            c.fill.solid(); c.fill.fore_color.rgb = NAVY
        else:
            pr.runs[0].font.color.rgb = GRAY
bx = s.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(12), Inches(0.8))
_set(bx.text_frame.paragraphs[0],
     "Giá trị trả về: NF_ACCEPT (cho qua)  ·  NF_DROP (loại bỏ gói)",
     16, NAVY, bold=True)

# =====================================================================
# 9. KERNEL SOCKET API
# =====================================================================
s = new_slide("Lập trình socket trong nhân", kicker="CƠ SỞ LÝ THUYẾT")
bullets(s, [
    "Nhân cung cấp API socket song song với user-space:",
], left=Inches(0.6), width=Inches(12), size=18, gap=8)
code_box(s,
    "sock_create_kern()   /* tạo socket trong nhân */\n"
    "kernel_bind()        /* ~ bind()   */\n"
    "kernel_listen()      /* ~ listen() */\n"
    "kernel_accept()      /* ~ accept() */\n"
    "kernel_sendmsg()     /* gửi dữ liệu — struct msghdr + kvec */\n"
    "kernel_recvmsg()     /* nhận dữ liệu */",
    top=Inches(2.4), height=Inches(3.2), size=15)

# =====================================================================
# 10. BÀI 1 - CODE
# =====================================================================
s = new_slide("Phần 1 — Giám sát gói tin (net_monitor)", kicker="TRIỂN KHAI")
purpose(s, "Nền tảng của mọi công cụ giám sát/IDS — thấy toàn bộ lưu lượng "
           "qua máy ở mức nhân (kể cả localhost, điều công cụ user-space khó làm).")
code_box(s,
    "static unsigned int hook_func(void *priv, struct sk_buff *skb,\n"
    "                              const struct nf_hook_state *state)\n"
    "{\n"
    "    struct iphdr *iph = ip_hdr(skb);\n"
    "    if (!iph) return NF_ACCEPT;\n"
    "    if (iph->protocol == IPPROTO_TCP) {\n"
    "        struct tcphdr *t = tcp_hdr(skb);\n"
    "        pr_info(\"TCP %pI4:%u -> %pI4:%u [%s%s]\\n\",\n"
    "            &iph->saddr, ntohs(t->source),\n"
    "            &iph->daddr, ntohs(t->dest),\n"
    "            t->syn?\"SYN \":\"\", t->ack?\"ACK\":\"\");\n"
    "    }\n"
    "    return NF_ACCEPT;   /* chỉ giám sát */\n"
    "}",
    top=Inches(2.25), height=Inches(4.35), width=Inches(7.2), size=12.5)
img_ph(s, "Ảnh chụp đăng ký hook trong hàm init "
          "(nfho.hooknum = NF_INET_PRE_ROUTING; nf_register_net_hook...).",
       left=Inches(8.1), width=Inches(4.7), top=Inches(2.25), height=Inches(4.35))

# =====================================================================
# 11. BÀI 1 - KẾT QUẢ + PHÂN TÍCH
# =====================================================================
s = new_slide("Phần 1 — Kết quả & phân tích log", kicker="THỰC NGHIỆM")
code_box(s,
    "net_monitor: TCP 192.168.0.102:53270 -> 192.168.0.108:22 [ACK ]\n"
    "net_monitor: ICMP 8.8.8.8 -> 192.168.0.105\n"
    "net_monitor: UDP 127.0.0.1:35681 -> 127.0.0.53:53\n"
    "net_monitor: UDP 8.8.8.8:53 -> 192.168.0.105:55135",
    top=Inches(1.45), height=Inches(1.7), width=Inches(12.1), size=13)
bullets(s, [
    "TCP :22 → phiên SSH đang điều khiển VM",
    "UDP → 127.0.0.53:53 → truy vấn DNS qua systemd-resolved",
    "UDP từ 8.8.8.8:53 → phản hồi DNS của Google",
    "ICMP từ 8.8.8.8 → gói trả lời ping",
    "→ Bắt được cả lưu lượng localhost lẫn ra ngoài.",
], top=Inches(3.4), size=17, gap=9)

# =====================================================================
# 12. BÀI 2 - CODE (ioctl + LOCAL_OUT)
# =====================================================================
s = new_slide("Phần 2 — Lọc gói theo IP đích (ioctl)", kicker="TRIỂN KHAI")
purpose(s, "Nguyên lý của firewall (iptables) — chủ động CHẶN lưu lượng theo "
           "chính sách, cấu hình động lúc chạy từ user-space.")
code_box(s,
    "/* hook đặt tại NF_INET_LOCAL_OUT */\n"
    "if (iph && iph->daddr == blocked_addr) {\n"
    "    pr_info(\"DROP goi den %pI4\\n\", &iph->daddr);\n"
    "    return NF_DROP;\n"
    "}\n"
    "return NF_ACCEPT;\n"
    "\n"
    "/* ioctl nhận IP cần chặn từ user-space */\n"
    "case FILTER_SET_ADDR:\n"
    "    copy_from_user(&addr, (void __user*)arg, sizeof(addr));\n"
    "    blocked_addr = addr;",
    top=Inches(2.25), height=Inches(3.85), size=13.5,
    caption="Mấu chốt: hook ở LOCAL_OUT để bắt gói ĐI RA (ping), không phải PRE_ROUTING.")

# =====================================================================
# 13. BÀI 2 - KẾT QUẢ
# =====================================================================
s = new_slide("Phần 2 — Kết quả chặn IP", kicker="THỰC NGHIỆM")
bullets(s, [
    "Quy trình demo:",
    ("", "sudo insmod net_filter.ko"),
    ("", "sudo ./filter_ctl set 8.8.8.8"),
    ("", "ping 8.8.8.8  →  100% packet loss"),
    ("", "sudo ./filter_ctl clear  →  ping hoạt động lại"),
], left=Inches(0.6), width=Inches(6.4), size=17, gap=11)
img_ph(s, "Ảnh terminal: ping 8.8.8.8 bị 100% packet loss SAU khi set, "
          "kèm log 'DROP goi den 8.8.8.8' trong dmesg.",
       left=Inches(7.3), width=Inches(5.4), top=Inches(1.5), height=Inches(4.6))

# =====================================================================
# 14. BÀI 3+4 - CODE
# =====================================================================
s = new_slide("Phần 3+4 — TCP listen/accept trong nhân", kicker="TRIỂN KHAI")
purpose(s, "Dịch vụ mạng chạy thẳng trong nhân (vd VPN, honeypot, giám sát) — "
           "lắng nghe & xử lý kết nối không cần tiến trình user-space.")
code_box(s,
    "sock_create_kern(&init_net, AF_INET, SOCK_STREAM,\n"
    "                 IPPROTO_TCP, &listen_sock);\n"
    "addr.sin_port = htons(60000);\n"
    "kernel_bind(listen_sock, (struct sockaddr*)&addr, sizeof(addr));\n"
    "kernel_listen(listen_sock, 5);\n"
    "\n"
    "while (!kthread_should_stop()) {\n"
    "    ret = kernel_accept(listen_sock, &conn, O_NONBLOCK);\n"
    "    if (ret == -EAGAIN) { msleep(200); continue; }\n"
    "    if (ret == 0) {\n"
    "        kernel_recvmsg(conn, &msg, &vec, 1, BUFSZ-1, 0);\n"
    "        sock_release(conn);\n"
    "    }\n"
    "}",
    top=Inches(2.25), height=Inches(4.35), size=12.5)

# =====================================================================
# 15. BÀI 3+4 - KẾT QUẢ
# =====================================================================
s = new_slide("Phần 3+4 — Kết quả", kicker="THỰC NGHIỆM")
bullets(s, [
    "accept không chặn (O_NONBLOCK + msleep) → rmmod gỡ sạch.",
    "Test từ user-space:",
    ("", 'echo "hello kernel" | nc 127.0.0.1 60000'),
    "dmesg in ra dữ liệu nhận được → socket nhân hoạt động.",
], left=Inches(0.6), width=Inches(6.4), size=17, gap=11)
img_ph(s, "Ảnh 2 terminal: bên trái gửi bằng nc tới cổng 60000, "
          "bên phải dmesg hiển thị 'nhan N byte: hello kernel'.",
       left=Inches(7.3), width=Inches(5.4), top=Inches(1.5), height=Inches(4.6))

# =====================================================================
# 16. BÀI 5 - CODE
# =====================================================================
s = new_slide("Phần 5 — Gửi UDP từ nhân (ksock_udp)", kicker="TRIỂN KHAI")
purpose(s, "Nhân chủ động GỬI dữ liệu ra mạng — vd đẩy cảnh báo/log/heartbeat "
           "tới máy giám sát ngay từ trong nhân.")
code_box(s,
    "sock_create_kern(&init_net, AF_INET, SOCK_DGRAM,\n"
    "                 IPPROTO_UDP, &sock);\n"
    "addr.sin_port        = htons(dport);\n"
    "addr.sin_addr.s_addr = in_aton(dip);\n"
    "\n"
    "mhdr.msg_name = &addr; mhdr.msg_namelen = sizeof(addr);\n"
    "vec.iov_base  = msg;   vec.iov_len = strlen(msg);\n"
    "\n"
    "kernel_sendmsg(sock, &mhdr, &vec, 1, vec.iov_len);\n"
    "sock_release(sock);",
    top=Inches(2.25), height=Inches(3.5), width=Inches(7.2), size=13)
img_ph(s, "Ảnh: lệnh insmod truyền tham số "
          "(dip, dport, 'msg=\"...\"') — chú ý cách quote chuỗi có khoảng trắng.",
       left=Inches(8.1), width=Inches(4.7), top=Inches(2.25), height=Inches(3.5))

# =====================================================================
# 17. BÀI 5 - KẾT QUẢ
# =====================================================================
s = new_slide("Phần 5 — Kết quả", kicker="THỰC NGHIỆM")
bullets(s, [
    "Mở listener trước, rồi nạp mô-đun:",
], left=Inches(0.6), width=Inches(12), size=17, gap=8)
code_box(s,
    "# Terminal 1:\n"
    "nc -u -l 5005\n"
    "# Terminal 2:\n"
    "sudo insmod ksock_udp.ko dip=\"127.0.0.1\" dport=5005 \\\n"
    "     'msg=\"Xin chao tu kernel\"'",
    top=Inches(2.3), height=Inches(1.9), width=Inches(7.2), size=13.5)
bullets(s, [
    "Lưu ý: msg có khoảng trắng → bọc nháy đơn quanh cả key=\"value\".",
    "Listener nhận đúng chuỗi gửi từ nhân.",
], top=Inches(4.5), left=Inches(0.6), width=Inches(7.2), size=15, gap=8)
img_ph(s, "Ảnh: terminal nc -u -l 5005 nhận được 'Xin chao tu kernel'.",
       left=Inches(8.1), width=Inches(4.7), top=Inches(2.3), height=Inches(4.1))

# =====================================================================
# 19. KẾT QUẢ TỔNG HỢP
# =====================================================================
s = new_slide("Kết quả tổng hợp", kicker="ĐÁNH GIÁ")
rows = [
    ("Phần", "Yêu cầu", "Kết quả"),
    ("1", "Giám sát / hiển thị gói tin", "Đạt"),
    ("2", "Lọc gói theo IP đích (ioctl)", "Đạt"),
    ("3", "TCP listen trong nhân", "Đạt"),
    ("4", "Accept kết nối trong nhân", "Đạt"),
    ("5", "Gửi UDP từ nhân", "Đạt"),
]
t = s.shapes.add_table(len(rows), 3, Inches(1.4), Inches(1.7),
                       Inches(10.5), Inches(3.8)).table
t.columns[0].width = Inches(1.2)
t.columns[1].width = Inches(7.0)
t.columns[2].width = Inches(2.3)
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        c = t.cell(i, j); c.text = val
        pr = c.text_frame.paragraphs[0]; r = pr.runs[0]
        r.font.size = Pt(16); r.font.name = SANS
        if i == 0:
            r.font.bold = True; r.font.color.rgb = WHITE
            c.fill.solid(); c.fill.fore_color.rgb = NAVY
        else:
            r.font.color.rgb = GRAY
            if j == 2:
                r.font.bold = True; r.font.color.rgb = RGBColor(0x1E, 0x8A, 0x3C)
bx = s.shapes.add_textbox(Inches(1.4), Inches(5.8), Inches(10.5), Inches(0.8))
_set(bx.text_frame.paragraphs[0],
     "Các mô-đun nạp/gỡ ổn định, không gây lỗi nhân.",
     16, NAVY, bold=True, align=PP_ALIGN.CENTER)

# =====================================================================
# 20. KẾT LUẬN
# =====================================================================
s = new_slide("Kết luận & hướng phát triển", kicker="KẾT LUẬN")
bullets(s, [
    "Hiểu kiến trúc nhân Linux: kernel space, kernel module, hệ thống mạng.",
    "Nắm được network stack, Netfilter, sk_buff & socket trong nhân.",
    "Triển khai thành công 5 phần kernel module trên Ubuntu kernel 5.15.",
    "Hướng phát triển:",
    ("", "Lọc theo cả cổng & giao thức, hỗ trợ IPv6"),
    ("", "Thống kê lưu lượng realtime"),
    ("", "Ghi log ra file/hệ thống giám sát tập trung"),
], size=18, gap=12)

# =====================================================================
# 21. CẢM ƠN
# =====================================================================
s = new_slide(count=False)
add_box(s, 0, 0, SW, SH, fill=NAVY)
tb = s.shapes.add_textbox(Inches(1), Inches(2.7), SW - Inches(2), Inches(1.5))
_set(tb.text_frame.paragraphs[0], "CẢM ƠN THẦY/CÔ & CÁC BẠN", 36, WHITE,
     bold=True, align=PP_ALIGN.CENTER)
tb = s.shapes.add_textbox(Inches(1), Inches(4.2), SW - Inches(2), Inches(1))
_set(tb.text_frame.paragraphs[0], "ĐÃ LẮNG NGHE — Q & A", 20,
     RGBColor(0x8E, 0xC6, 0xFF), align=PP_ALIGN.CENTER)

prs.save(OUT)
print("Da tao:", OUT, "|", len(prs.slides._sldIdLst), "slides")
